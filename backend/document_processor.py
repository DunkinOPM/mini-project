"""
Document Processing Module
- Extracts text from .pptx and .pdf files
- Returns structured slide/page data
"""

import re
from pathlib import Path
from typing import List, Dict, Any


class DocumentProcessor:
    def process(self, doc_path: str) -> List[Dict[str, Any]]:
        path = Path(doc_path)
        ext = path.suffix.lower()
        if ext == ".pptx":
            return self._process_pptx(doc_path)
        elif ext == ".pdf":
            return self._process_pdf(doc_path)
        else:
            raise ValueError(f"Unsupported document format: {ext}")

    # ──────────────────────────────────────────────
    # PPT
    # ──────────────────────────────────────────────
    def _process_pptx(self, path: str) -> List[Dict[str, Any]]:
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            print("[DocumentProcessor] python-pptx not installed. Using mock slides.")
            return self._mock_slides()

        prs = Presentation(path)
        slides = []
        for idx, slide in enumerate(prs.slides):
            title = ""
            bullets = []
            full_text = ""

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for i, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if not text:
                        continue
                    if i == 0 and not title:
                        title = text
                    else:
                        bullets.append(text)
                    full_text += text + " "

            slides.append({
                "id": idx,
                "slide_number": idx + 1,
                "title": title or f"Slide {idx + 1}",
                "bullets": bullets,
                "text": full_text.strip(),
                "sections": self._split_sections(full_text),
            })
        return slides

    # ──────────────────────────────────────────────
    # PDF
    # ──────────────────────────────────────────────
    def _process_pdf(self, path: str) -> List[Dict[str, Any]]:
        # Try PyMuPDF first
        try:
            import fitz  # PyMuPDF
            return self._pdf_with_fitz(path)
        except ImportError:
            pass

        # Fallback to pdfplumber
        try:
            import pdfplumber
            return self._pdf_with_pdfplumber(path)
        except ImportError:
            pass

        print("[DocumentProcessor] No PDF lib found. Using mock slides.")
        return self._mock_slides()

    def _pdf_with_fitz(self, path: str) -> List[Dict[str, Any]]:
        import fitz
        doc = fitz.open(path)
        slides = []
        for idx, page in enumerate(doc):
            text = page.get_text()
            lines = [l.strip() for l in text.split("\n") if l.strip()]
            title = lines[0] if lines else f"Page {idx + 1}"
            bullets = lines[1:] if len(lines) > 1 else []
            slides.append({
                "id": idx,
                "slide_number": idx + 1,
                "title": title,
                "bullets": bullets,
                "text": text.strip(),
                "sections": self._split_sections(text),
            })
        return slides

    def _pdf_with_pdfplumber(self, path: str) -> List[Dict[str, Any]]:
        import pdfplumber
        slides = []
        with pdfplumber.open(path) as pdf:
            for idx, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                lines = [l.strip() for l in text.split("\n") if l.strip()]
                title = lines[0] if lines else f"Page {idx + 1}"
                bullets = lines[1:] if len(lines) > 1 else []
                slides.append({
                    "id": idx,
                    "slide_number": idx + 1,
                    "title": title,
                    "bullets": bullets,
                    "text": text.strip(),
                    "sections": self._split_sections(text),
                })
        return slides

    # ──────────────────────────────────────────────
    # Helpers
    # ──────────────────────────────────────────────
    def _split_sections(self, text: str) -> List[str]:
        """Split text into paragraph-level sections."""
        paras = re.split(r"\n{2,}|\.\s{2,}", text)
        return [p.strip() for p in paras if p.strip() and len(p.strip()) > 20]

    def _mock_slides(self) -> List[Dict[str, Any]]:
        """Fallback mock slides."""
        data = [
            ("Introduction to Machine Learning", [
                "Definition: ML is a subset of AI",
                "Types: Supervised, Unsupervised, Reinforcement",
                "Applications: image recognition, NLP, robotics",
            ]),
            ("Supervised Learning", [
                "Training on labeled data",
                "Classification and Regression tasks",
                "Examples: Decision Trees, SVM, Neural Networks",
            ]),
            ("Decision Trees", [
                "Tree-like model of decisions",
                "Splits data on feature values",
                "Handles both categorical and numerical data",
            ]),
            ("Neural Networks", [
                "Inspired by the human brain",
                "Layers: Input, Hidden, Output",
                "Activation functions: ReLU, Sigmoid, Softmax",
            ]),
            ("Gradient Descent", [
                "Optimization algorithm",
                "Minimizes loss function",
                "Learning rate controls step size",
            ]),
            ("Overfitting & Regularization", [
                "Overfitting: model memorizes training data",
                "Regularization: L1, L2 penalties",
                "Dropout technique in neural networks",
            ]),
            ("Cross-Validation", [
                "K-fold cross-validation",
                "Train/validation/test splits",
                "Evaluating generalization performance",
            ]),
            ("Convolutional Neural Networks", [
                "Specialized for image data",
                "Convolution, pooling, fully-connected layers",
                "Applications: object detection, classification",
            ]),
            ("Transfer Learning", [
                "Reuse pre-trained models",
                "Fine-tuning for specific tasks",
                "Reduces data and compute requirements",
            ]),
        ]
        slides = []
        for idx, (title, bullets) in enumerate(data):
            text = title + ". " + ". ".join(bullets)
            slides.append({
                "id": idx,
                "slide_number": idx + 1,
                "title": title,
                "bullets": bullets,
                "text": text,
                "sections": [text],
            })
        return slides
